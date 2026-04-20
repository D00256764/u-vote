#!/usr/bin/env python3
"""
U-Vote Stage 2 — Presentation Generator
Style: IoT Platform deck — big bold titles, short bullets, speaker fills in detail
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

OFF_WHITE    = RGBColor(0xF8, 0xF8, 0xFF)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x0D, 0x0D, 0x1A)
PURPLE_DARK  = RGBColor(0x3B, 0x1F, 0x9E)
PURPLE_MID   = RGBColor(0x6D, 0x50, 0xE0)
PURPLE_LIGHT = RGBColor(0xC4, 0xB5, 0xFD)
GREY_LIGHT   = RGBColor(0xEE, 0xED, 0xF8)
GREY_MID     = RGBColor(0x44, 0x44, 0x60)
GREEN        = RGBColor(0x16, 0xA3, 0x4A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]
NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def bg(slide):
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = OFF_WHITE
    s.line.fill.background()


def blob(slide, cx, cy, size, color=PURPLE_LIGHT, alpha=40):
    s = slide.shapes.add_shape(9, cx - size/2, cy - size/2, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    sf = s._element.find('.//{%s}solidFill' % NS)
    if sf is not None:
        clr = sf.find('{%s}srgbClr' % NS)
        if clr is not None:
            a = etree.SubElement(clr, '{%s}alpha' % NS)
            a.set('val', str(alpha * 1000))


def rect(slide, x, y, w, h, fill, border=None, lw=0.75):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def tb(slide, text, x, y, w, h,
       size=18, bold=False, italic=False,
       color=BLACK, align=PP_ALIGN.LEFT, font="Calibri"):
    t = slide.shapes.add_textbox(x, y, w, h)
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = font
    return t


def kv_bullets(slide, items, x, y, w, size=20, gap=Pt(14)):
    txb = slide.shapes.add_textbox(x, y, w, Inches(6))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if not first: p.space_before = gap
        first = False
        p.alignment = PP_ALIGN.LEFT
        if isinstance(item, tuple):
            term, desc = item
            r1 = p.add_run()
            r1.text = f"\u2022   {term}: "
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = BLACK; r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(size); r2.font.bold = False
            r2.font.color.rgb = GREY_MID; r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = f"\u2022   {item}" if item else ""
            r.font.size = Pt(size); r.font.color.rgb = GREY_MID
            r.font.name = "Calibri"


def job_table(slide, rows, x, y, w, col1w=Inches(2.1)):
    """Render a two-column job table (Job label | description)."""
    col2w = w - col1w
    row_h = Inches(0.52)
    for i, (job, desc) in enumerate(rows):
        ry = y + i * row_h
        fill = PURPLE_DARK if i == 0 else (GREY_LIGHT if i % 2 == 1 else WHITE)
        tc  = WHITE if i == 0 else BLACK
        rect(slide, x, ry, col1w, row_h, fill,
             border=PURPLE_DARK if i == 0 else PURPLE_LIGHT)
        rect(slide, x + col1w, ry, col2w, row_h, fill,
             border=PURPLE_DARK if i == 0 else PURPLE_LIGHT)
        tb(slide, job,
           x + Inches(0.08), ry + Inches(0.08), col1w - Inches(0.1), row_h - Inches(0.1),
           size=12, bold=(i == 0), color=tc)
        tb(slide, desc,
           x + col1w + Inches(0.1), ry + Inches(0.08), col2w - Inches(0.15), row_h - Inches(0.1),
           size=12, color=tc)


def box_node(slide, text, x, y, w, h,
             fill=GREY_LIGHT, border=PURPLE_MID,
             size=12, bold=False, color=BLACK):
    rect(slide, x, y, w, h, fill, border)
    tb(slide, text, x, y + (h - Inches(0.3)) / 2, w, Inches(0.3),
       size=size, bold=bold, color=color, align=PP_ALIGN.CENTER)


def arrow_h(slide, x1, ym, x2):
    c = slide.shapes.add_connector(1, x1, ym, x2, ym)
    c.line.color.rgb = GREY_MID; c.line.width = Pt(1.5)


def arrow_v(slide, xm, y1, y2):
    c = slide.shapes.add_connector(1, xm, y1, xm, y2)
    c.line.color.rgb = GREY_MID; c.line.width = Pt(1.5)


def big_title(slide, line1, line2=None):
    bg(slide)
    blob(slide, Inches(1.5),  Inches(5.5), Inches(5),   PURPLE_LIGHT, 38)
    blob(slide, Inches(11.5), Inches(1.5), Inches(3.5), PURPLE_MID,   28)
    tb(slide, line1, Inches(0.6), Inches(0.4), Inches(12), Inches(2.2),
       size=72, bold=True, color=BLACK)
    if line2:
        tb(slide, line2, Inches(0.6), Inches(2.5), Inches(12), Inches(1.5),
           size=72, bold=True, color=BLACK)


def content(slide, title, subtitle=None):
    bg(slide)
    blob(slide, Inches(11.5), Inches(5.5), Inches(4), PURPLE_LIGHT, 35)
    tb(slide, title, Inches(0.6), Inches(0.3), Inches(12), Inches(1.0),
       size=44, bold=True, color=BLACK)
    if subtitle:
        tb(slide, subtitle, Inches(0.6), Inches(1.25), Inches(12), Inches(0.4),
           size=16, color=GREY_MID)
    rect(slide, Inches(0.6), Inches(1.6), Inches(12.35), Inches(0.03), PURPLE_MID)


def diagram_placeholder(slide, label="Insert LucidChart diagram here"):
    rect(slide, Inches(0.9), Inches(1.82), Inches(11.55), Inches(5.35),
         GREY_LIGHT, border=PURPLE_LIGHT, lw=1.5)
    tb(slide, label, Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.6),
       size=20, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)


def stat_cards(slide, cards, x, y):
    """Render a row of stat cards: (number, label, sub)."""
    CW = Inches(2.8); CH = Inches(1.1); gap = Inches(0.15)
    for i, (num, label, sub) in enumerate(cards):
        cx = x + i * (CW + gap)
        rect(slide, cx, y, CW, CH, GREY_LIGHT, border=PURPLE_LIGHT)
        tb(slide, num, cx + Inches(0.12), y + Inches(0.08),
           Inches(1.1), Inches(0.5), size=30, bold=True, color=PURPLE_DARK)
        tb(slide, label, cx + Inches(1.25), y + Inches(0.1),
           Inches(1.45), Inches(0.38), size=12, bold=True, color=BLACK)
        tb(slide, sub, cx + Inches(1.25), y + Inches(0.52),
           Inches(1.45), Inches(0.42), size=10, color=GREY_MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
blob(s, Inches(1.5),  Inches(5.5), Inches(5.5), PURPLE_LIGHT, 38)
blob(s, Inches(11.5), Inches(1.5), Inches(4.5), PURPLE_MID,   28)
blob(s, Inches(7),    Inches(6.5), Inches(3),   PURPLE_LIGHT, 30)
rect(s, Inches(0.55), Inches(1.3), Inches(0.1), Inches(4.5), PURPLE_DARK)
tb(s, "U-VOTE",
   Inches(0.8), Inches(1.25), Inches(11), Inches(1.7),
   size=88, bold=True, color=BLACK)
tb(s, "SECURE ONLINE\nVOTING SYSTEM",
   Inches(0.8), Inches(2.9), Inches(11), Inches(2.2),
   size=46, bold=True, color=BLACK)
tb(s, "MVP Release \u2014 From Development to Production",
   Inches(0.8), Inches(5.2), Inches(10), Inches(0.5),
   size=20, color=GREY_MID)
tb(s, "Luke Doyle (D00255656)   \u00b7   Hafsa Moin (D00256764)",
   Inches(0.8), Inches(5.75), Inches(10), Inches(0.45),
   size=16, color=GREY_MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Introduction")

tb(s, "Goal:",
   Inches(0.6), Inches(1.82), Inches(1.0), Inches(0.42),
   size=20, bold=True, color=BLACK)
tb(s, "A production-grade, secure online voting platform built for small-scale organisations \u2014 designed for verifiability, anonymity, and auditability.",
   Inches(1.65), Inches(1.82), Inches(11.0), Inches(0.42),
   size=20, color=GREY_MID)

tb(s, "Summary:",
   Inches(0.6), Inches(2.62), Inches(2.0), Inches(0.42),
   size=20, bold=True, color=BLACK)

kv_bullets(s, [
    "Fully automated CI/CD pipeline with 9 quality-gated stages from lint through to deployment",
    "Production-grade Kubernetes cluster with Calico zero-trust network policies",
    "Complete working MVP \u2014 phone OTP MFA, automatic election scheduling, cryptographic audit trail",
    "Live observability \u2014 Prometheus, Grafana, Fluent Bit, Elasticsearch, and Kibana",
], Inches(0.6), Inches(3.15), Inches(12.0), size=20, gap=Pt(12))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS U-VOTE?
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
big_title(s, "What is", "U-Vote?")

kv_bullets(s, [
    ("auth-service", "JWT authentication, phone OTP MFA, blind ballot token issuance"),
    ("election-service", "Election lifecycle management, scheduled open/close, voter management"),
    ("voting-service", "Voter-facing ballot submission and receipt verification"),
    ("results-service", "Vote tallying, winner calculation, and immutable audit log"),
    ("admin-service", "Voter CSV import, token generation, and GDPR erasure"),
    ("frontend-service", "Organiser dashboard and election management UI"),
], Inches(0.6), Inches(4.55), Inches(12.0), size=20, gap=Pt(11))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "System Architecture",
        "Browser \u2192 NGINX Ingress \u2192 6 Services \u2192 PostgreSQL  \u00b7  All services \u2192 Prometheus + ELK")
diagram_placeholder(s, "Insert architecture diagram (LucidChart)")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CONTINUOUS INTEGRATION PIPELINE (big title)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
big_title(s, "Continuous", "Integration Pipeline")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GITHUB ACTIONS — JOB TABLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "GitHub Actions",
        "9 jobs \u2014 triggered on every push to every branch")

job_table(s, [
    ("Job",         "Description"),
    ("1  Lint",     "ruff check across all Python source files \u2014 blocks all downstream jobs on failure"),
    ("2  Unit Tests","pytest matrix: 6 services in parallel \u00b7 121 tests \u00b7 JUnit XML artefacts uploaded per service \u00b7 fail-fast: false"),
    ("3  Shared Tests","tests/test_shared.py & tests/test_logging_config.py"),
    ("4  Integration","Docker Compose full stack spin-up \u00b7 NGINX gateway \u00b7 HTTP flows: register \u2192 login \u2192 dashboard \u00b7 tears down with docker compose down -v"),
    ("5  Docker Build","docker compose build \u2014 verifies all 6 Dockerfiles build cleanly"),
    ("6  Trivy CVE", "Matrix scan: all 6 images \u00b7 CRITICAL and HIGH severity CVEs \u00b7 SARIF uploaded to GitHub Security tab"),
    ("7  Trivy FS",  "Scans source code for hardcoded secrets, exposed credentials, and infrastructure misconfigurations"),
    ("8  Push GHCR", "Builds and pushes all 6 images tagged with commit SHA \u2014 main and mvp_tests branches only"),
    ("9  Manifests", "Stamps new image SHA into Kubernetes deployment YAMLs and commits back \u2014 ArgoCD detects and syncs"),
],
    x=Inches(0.6), y=Inches(1.72), w=Inches(12.35))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PIPELINE FLOW DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Pipeline Flow",
        "Lint \u2192 Tests \u2192 Build \u2192 Trivy Scan \u2192 Push Registry \u2192 Update Manifests \u2192 ArgoCD Sync")
diagram_placeholder(s, "Insert pipeline flow diagram (LucidChart)")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — KUBERNETES PLATFORM (big title)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
big_title(s, "Kubernetes Platform", "\u2014 KIND")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CLUSTER TOPOLOGY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Cluster Topology")

tb(s, "Cluster:",
   Inches(0.6), Inches(1.82), Inches(1.4), Inches(0.42),
   size=19, bold=True, color=BLACK)
kv_bullets(s, [
    "1 control-plane node  \u00b7  2 worker nodes",
    "Pod subnet: 192.168.0.0/16  \u00b7  CNI: Calico",
    "2 replicas per service  \u00b7  rolling updates  \u00b7  liveness + readiness probes",
    "All credentials in Kubernetes Secrets \u2014 never in code",
], Inches(0.6), Inches(2.35), Inches(6.0), size=18, gap=Pt(9))

rect(s, Inches(6.75), Inches(1.75), Inches(0.02), Inches(5.4), PURPLE_LIGHT)

tb(s, "Namespaces:",
   Inches(6.9), Inches(1.82), Inches(5.9), Inches(0.42),
   size=19, bold=True, color=BLACK)

# Namespace status table
ns_rows = [
    ("uvote-dev",  "Development and feature testing",      GREEN),
    ("uvote-test", "Integration testing \u2014 mirrors prod",     GREEN),
    ("uvote-prod", "Production-grade deployment",           PURPLE_MID),
]
for i, (ns, desc, col) in enumerate(ns_rows):
    ry = Inches(2.38 + i * 0.72)
    rect(s, Inches(6.9), ry, Inches(5.9), Inches(0.58), GREY_LIGHT, border=PURPLE_LIGHT)
    tb(s, ns,   Inches(7.05), ry + Inches(0.08), Inches(1.6),  Inches(0.42), size=14, bold=True, color=col)
    tb(s, desc, Inches(8.75), ry + Inches(0.08), Inches(3.85), Inches(0.42), size=13, color=GREY_MID)

diagram_placeholder_small = lambda sl: (
    rect(sl, Inches(0.6), Inches(5.4), Inches(5.9), Inches(0.4), GREY_LIGHT, border=PURPLE_LIGHT),
    tb(sl, "Insert cluster topology diagram (LucidChart)",
       Inches(0.65), Inches(5.42), Inches(5.8), Inches(0.35),
       size=11, italic=True, color=GREY_MID)
)
diagram_placeholder_small(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — NETWORK POLICIES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Network Policies",
        "Calico zero-trust \u2014 default deny, explicit allow")

kv_bullets(s, [
    ("Policy 00", "Default deny \u2014 all ingress and egress blocked for every pod"),
    ("Policy 01", "Allow DNS \u2014 all pods can resolve internal service names via CoreDNS"),
    ("Policy 02", "Allow service \u2192 PostgreSQL \u2014 per-service scoped database access only"),
    ("Policy 03", "Allow ingress controller \u2192 services \u2014 only NGINX can reach backend pods"),
    ("Policy 04", "Allow audit service access \u2014 inter-service calls for audit logging"),
    ("Policy 05", "Allow MailHog SMTP \u2014 email capture in dev and test environments only"),
    ("Policy 06", "Allow Prometheus scraping \u2014 monitoring namespace can reach metrics endpoints"),
], Inches(0.6), Inches(1.82), Inches(12.0), size=19, gap=Pt(10))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — OBSERVABILITY STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Observability Stack",
        "Metrics \u00b7 Logs \u00b7 Live dashboards")

tb(s, "Metrics:",
   Inches(0.6), Inches(1.82), Inches(1.3), Inches(0.42),
   size=18, bold=True, color=PURPLE_DARK)
kv_bullets(s, [
    ("Prometheus", "Scrapes /metrics from all 6 services every 15 seconds"),
    ("Grafana", "SLI dashboard \u2014 request rate, p95 latency, 5xx error rate, availability"),
], Inches(0.6), Inches(2.35), Inches(5.9), size=17, gap=Pt(10))

rect(s, Inches(6.75), Inches(1.75), Inches(0.02), Inches(5.4), PURPLE_LIGHT)

tb(s, "Logs:",
   Inches(6.9), Inches(1.82), Inches(1.0), Inches(0.42),
   size=18, bold=True, color=PURPLE_DARK)
kv_bullets(s, [
    ("Fluent Bit", "Collects container logs from all pods as a DaemonSet on every node"),
    ("Elasticsearch", "Indexes all logs under the uvote-logs* index pattern"),
    ("Kibana", "Live dashboards \u2014 error volume, auth failure rate, vote submission rate, live log stream"),
], Inches(6.9), Inches(2.35), Inches(5.9), size=17, gap=Pt(10))

# SLO bar
rect(s, Inches(0.55), Inches(6.35), Inches(12.3), Inches(0.68), GREY_LIGHT)
tb(s, "SLO targets:   Availability \u2265 99.5%   \u00b7   P99 latency < 500ms   \u00b7   Error rate < 0.1%   \u00b7   No PII written to logs",
   Inches(0.72), Inches(6.4), Inches(12.0), Inches(0.58),
   size=14, bold=True, color=PURPLE_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DEPLOYMENT PROCESS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Deployment Process",
        "Full cluster bootstrapped and verified in a single automated sequence")

tb(s, "Bootstrap sequence:",
   Inches(0.6), Inches(1.82), Inches(5.9), Inches(0.42),
   size=18, bold=True, color=PURPLE_DARK)
kv_bullets(s, [
    ("Step 1 \u2014 Setup", "Virtual environment created \u2014 all dependencies installed for 6 services, platform scripts, and test suites"),
    ("Step 2a \u2014 Cluster", "KIND cluster \u2192 Calico CNI \u2192 namespaces \u2192 PostgreSQL \u2192 schema \u2192 roles \u2192 network policies \u2192 NGINX ingress"),
    ("Step 2b \u2014 Services", "Docker build (6 images) \u2192 kind load \u2192 Kubernetes Secrets \u2192 kubectl apply \u2192 pod readiness \u2192 network policy tests \u2192 health checks"),
    ("Step 2c \u2014 ELK", "Helm install Elasticsearch \u2192 Kibana \u2192 Fluent Bit (monitoring namespace) \u2192 Kibana dashboards auto-created"),
    ("Step 2d \u2014 Mail", "MailHog SMTP capture deployed to uvote-dev and uvote-test"),
], Inches(0.6), Inches(2.35), Inches(5.9), size=15, gap=Pt(8))

rect(s, Inches(6.75), Inches(1.75), Inches(0.02), Inches(5.4), PURPLE_LIGHT)

tb(s, "Result:",
   Inches(6.9), Inches(1.82), Inches(5.9), Inches(0.42),
   size=18, bold=True, color=PURPLE_DARK)

stat_cards(s, [
    ("~10m", "total time", "single command"),
    ("13/13", "pods running", "all healthy"),
], Inches(6.9), Inches(2.42))

stat_cards(s, [
    ("6/6",  "health checks", "all services"),
    ("7/7",  "network tests", "zero-trust verified"),
], Inches(6.9), Inches(3.72))

tb(s, "ArgoCD:",
   Inches(6.9), Inches(4.95), Inches(5.9), Inches(0.42),
   size=16, bold=True, color=PURPLE_DARK)
kv_bullets(s, [
    "Watches the Git repo for manifest changes",
    "Syncs cluster automatically on every CI push to main",
    "Rollback = git revert \u2014 ArgoCD re-syncs to previous version",
], Inches(6.9), Inches(5.42), Inches(5.9), size=15, gap=Pt(7))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — TESTING RESULTS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Testing Results")

stat_cards(s, [
    ("121",  "unit tests",    "6 services, all passing"),
    ("11/11","database tests","schema \u00b7 triggers \u00b7 roles"),
    ("28/28","API integration","full voter journey"),
    ("13/13","pods running",  "6 services \u00d7 2 + PostgreSQL"),
], Inches(0.6), Inches(1.82))

stat_cards(s, [
    ("6/6",  "health checks", "all services responding"),
    ("7/7",  "network tests", "zero-trust enforcement verified"),
], Inches(0.6), Inches(3.15))

tb(s, "Integration test coverage:",
   Inches(0.6), Inches(4.35), Inches(5.9), Inches(0.42),
   size=17, bold=True, color=PURPLE_DARK)
kv_bullets(s, [
    "Full Docker Compose stack spun up in CI \u2014 PostgreSQL, all 6 services, and NGINX gateway",
    "Tests cover: health endpoints \u00b7 organiser registration and login \u00b7 election creation \u00b7 CSV voter import \u00b7 vote submission \u00b7 results retrieval",
], Inches(0.6), Inches(4.85), Inches(5.9), size=15, gap=Pt(8))

rect(s, Inches(6.75), Inches(4.28), Inches(0.02), Inches(2.87), PURPLE_LIGHT)

tb(s, "Security scan coverage:",
   Inches(6.9), Inches(4.35), Inches(5.9), Inches(0.42),
   size=17, bold=True, color=PURPLE_DARK)
kv_bullets(s, [
    "Trivy scans all 6 images for CRITICAL and HIGH CVEs on every CI run",
    "Filesystem scan checks source code for hardcoded secrets and misconfigurations",
    "SARIF results uploaded to GitHub Security tab automatically",
], Inches(6.9), Inches(4.85), Inches(5.9), size=15, gap=Pt(8))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — VOTING JOURNEY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Voting Journey",
        "End-to-end \u2014 organiser setup through to published results")
diagram_placeholder(s, "Insert voting journey flow diagram (from Stage 1 LucidChart)")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SECURITY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Security")

kv_bullets(s, [
    ("Authentication", "JWT tokens \u00b7 bcrypt password hashing \u00b7 session management built into the auth service"),
    ("Phone OTP MFA", "6-digit code \u00b7 cryptographically secure generation \u00b7 hashed with pbkdf2_hmac (100k iterations) \u00b7 10-minute expiry \u00b7 max 3 attempts"),
    ("Ballot Anonymity", "Blind ballot tokens \u2014 voter identity permanently decoupled from vote content"),
    ("Vote Integrity", "Immutable ballots \u00b7 SHA-256 hash chain on every vote \u00b7 audit log protected by database trigger"),
    ("Double-Vote Prevention", "UNIQUE constraint on ballot token \u2014 resubmission rejected at the database level"),
    ("Least Privilege", "Per-service PostgreSQL users \u2014 each service can only access its own data"),
    ("Network Isolation", "Calico zero-trust \u00b7 default-deny \u00b7 explicit allow per service pair"),
], Inches(0.6), Inches(1.82), Inches(12.0), size=18, gap=Pt(9))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — LIVE MVP DEMO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
big_title(s, "Live MVP", "Demo")

kv_bullets(s, [
    "Organiser registers, logs in, creates an election with candidates and a scheduled open and close time",
    "Voters uploaded via CSV \u2014 election opens automatically, tokens emailed to each voter instantly on open",
    "Voter clicks unique link, completes phone OTP MFA, and casts ballot anonymously",
    "Results tallied and published automatically on close \u2014 full cryptographic audit trail available",
], Inches(0.6), Inches(4.55), Inches(12.0), size=22, gap=Pt(14))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — REQUIREMENTS FULFILLMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Requirements Fulfillment")

kv_bullets(s, [
    ("Secure Auth", "JWT tokens with bcrypt hashing \u2014 session management and login flow"),
    ("Phone OTP MFA", "6-digit codes delivered by email \u2014 hashed, expiring, brute-force protected"),
    ("Vote Integrity & Anonymity", "AES-256 encrypted votes \u00b7 SHA-256 hash chain \u00b7 blind tokens decouple voter from vote"),
    ("Auto Election Windows", "Background scheduler opens and closes elections at configured times automatically"),
    ("NGINX Gateway", "All traffic routed through NGINX Ingress Controller \u2014 single entry point"),
    ("CI/CD Pipeline", "9-stage GitHub Actions workflow \u2014 lint, test, scan, build, push, and deploy on every commit"),
    ("Observability", "Prometheus + Grafana SLI metrics \u00b7 Fluent Bit + Elasticsearch + Kibana log dashboards"),
    ("GDPR", "Hard-delete of all voter PII on request \u00b7 no PII in logs \u00b7 consent shown before ballot submission"),
], Inches(0.6), Inches(1.82), Inches(12.0), size=17, gap=Pt(7))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — NEXT STEPS & AWS ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content(s, "Next Steps & AWS Roadmap",
        "Same platform \u2014 same manifests \u2014 managed cloud infrastructure on AWS")

tb(s, "Phase 1 \u2014 AWS Deployment:",
   Inches(0.6), Inches(1.82), Inches(5.9), Inches(0.42),
   size=17, bold=True, color=PURPLE_DARK)
kv_bullets(s, [
    ("AWS EKS", "Managed Kubernetes cluster \u2014 same manifests, same ArgoCD config, no code changes required"),
    ("Terraform", "VPC, node groups, IAM roles, and ALB provisioned as code and version-controlled in Git"),
    ("RDS PostgreSQL", "Managed database with automated backups, multi-AZ failover, and point-in-time recovery"),
    ("TLS & DNS", "AWS Certificate Manager for HTTPS \u00b7 Route 53 for DNS routing to the load balancer"),
], Inches(0.6), Inches(2.35), Inches(5.9), size=16, gap=Pt(9))

rect(s, Inches(6.75), Inches(1.75), Inches(0.02), Inches(5.4), PURPLE_LIGHT)

tb(s, "Phase 2\u20135 \u2014 Beyond MVP:",
   Inches(6.9), Inches(1.82), Inches(5.9), Inches(0.42),
   size=17, bold=True, color=PURPLE_DARK)
kv_bullets(s, [
    ("Full CD", "ArgoCD promotes verified builds from test cluster to production automatically"),
    ("Multi-Region", "Secondary region with cross-region DB replication and Route 53 active/active failover"),
    ("Advanced Security", "WebAuthn hardware key support \u00b7 AWS HSM for key management \u00b7 third-party security audit"),
    ("Scale", "Horizontal pod autoscaling \u00b7 DB read replicas \u00b7 Redis cache \u00b7 CDN for static assets"),
], Inches(6.9), Inches(2.35), Inches(5.9), size=16, gap=Pt(9))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
blob(s, Inches(1),    Inches(1),   Inches(6),   PURPLE_LIGHT, 38)
blob(s, Inches(11.5), Inches(5.5), Inches(5),   PURPLE_MID,   28)
blob(s, Inches(6.5),  Inches(6),   Inches(3.5), PURPLE_LIGHT, 40)
tb(s, "THANK YOU",
   Inches(1.5), Inches(2.4), Inches(10.5), Inches(1.9),
   size=88, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
tb(s, "Luke Doyle (D00255656)   \u00b7   Hafsa Moin (D00256764)",
   Inches(1.5), Inches(4.4), Inches(10.5), Inches(0.55),
   size=24, color=GREY_MID, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "/Users/hafsa/Documents/u-vote/U-Vote-Stage2.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
